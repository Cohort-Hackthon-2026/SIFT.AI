"""Server-side export engine (BE2, plan §6 P2/P3).

Turns a chat transcript (question + grounded answer + citations) into a
downloadable memo/brief in PDF, DOCX, or PPTX. All three renderers consume one
intermediate :class:`ExportDocument` so the layout stays consistent across
formats. Heavy libraries are imported lazily and a missing one raises
:class:`ExportError` (the route maps it to HTTP 503) rather than crashing boot.
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field
from datetime import datetime, timezone
from typing import Any
from xml.sax.saxutils import escape

SUPPORTED_FORMATS = {"PDF", "DOCX", "PPTX"}

MIME_TYPES = {
    "PDF": "application/pdf",
    "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "PPTX": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
EXTENSIONS = {"PDF": "pdf", "DOCX": "docx", "PPTX": "pptx"}

_ROLE_LABELS = {"user": "Question", "assistant": "SIFT.AI Analysis", "system": "Context"}


class ExportError(Exception):
    """Raised when a format is unsupported or its renderer library is missing."""


@dataclass
class ExportBlock:
    heading: str
    paragraphs: list[str] = field(default_factory=list)
    citations: list[str] = field(default_factory=list)


@dataclass
class ExportDocument:
    title: str
    subtitle: str | None
    meta_lines: list[str]
    blocks: list[ExportBlock]


def utc_now() -> datetime:
    return datetime.now(timezone.utc)


# --------------------------------------------------------------------------- #
# Build the intermediate representation from a chat + messages
# --------------------------------------------------------------------------- #

def _format_citation(citation: Any) -> str:
    """Render one citation as a single readable line, tolerant of shape."""
    if isinstance(citation, str):
        return citation
    if not isinstance(citation, dict):
        return str(citation)
    name = citation.get("document_name") or citation.get("source") or citation.get("title") or "Source"
    page = citation.get("page") or citation.get("page_number")
    snippet = citation.get("snippet") or citation.get("text") or ""
    web_url = citation.get("url")
    parts = [str(name)]
    if page is not None:
        parts.append(f"p.{page}")
    line = " ".join(parts)
    if web_url:
        line = f"{line} — {web_url}"
    if snippet:
        snippet = snippet.strip().replace("\n", " ")
        if len(snippet) > 280:
            snippet = snippet[:277] + "…"
        line = f"{line}: “{snippet}”"
    return line


def _split_paragraphs(text: str) -> list[str]:
    if not text:
        return []
    # Collapse Windows newlines, split on blank lines, keep non-empty.
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    chunks = [c.strip() for c in normalized.split("\n\n")]
    return [c for c in chunks if c]


def build_chat_document(
    chat: dict[str, Any],
    messages: list[dict[str, Any]],
    *,
    generated_by: str | None = None,
    matter: dict[str, Any] | None = None,
    chambers: dict[str, Any] | None = None,
) -> ExportDocument:
    title = chat.get("title") or "SIFT.AI Research Memo"
    mode = (chat.get("mode") or "STRICT").upper()

    meta_lines = [f"Mode: {mode}", f"Generated: {utc_now().strftime('%Y-%m-%d %H:%M UTC')}"]
    if generated_by:
        meta_lines.append(f"Prepared by: {generated_by}")
    if chambers and chambers.get("name"):
        meta_lines.append(f"Chambers: {chambers['name']}")
    if matter and matter.get("title"):
        meta_lines.append(f"Matter: {matter['title']}")
        if matter.get("client_name"):
            meta_lines.append(f"Client: {matter['client_name']}")

    blocks: list[ExportBlock] = []
    for msg in messages:
        role = (msg.get("role") or "assistant").lower()
        heading = _ROLE_LABELS.get(role, role.title())
        paragraphs = _split_paragraphs(msg.get("content") or "")
        # Citations may be attached directly or (as the chat stream stores them)
        # split across metadata.internal_citations + metadata.external_citations.
        citations_raw = msg.get("citations")
        if citations_raw is None:
            meta = msg.get("metadata") or {}
            citations_raw = list(meta.get("internal_citations") or []) + list(
                meta.get("external_citations") or []
            )
        citations = [_format_citation(c) for c in citations_raw] if isinstance(citations_raw, list) else []
        if not paragraphs and not citations:
            continue
        blocks.append(ExportBlock(heading=heading, paragraphs=paragraphs, citations=citations))

    return ExportDocument(
        title=title,
        subtitle="AI-assisted legal research — verify all citations before filing.",
        meta_lines=meta_lines,
        blocks=blocks,
    )


# --------------------------------------------------------------------------- #
# Renderers
# --------------------------------------------------------------------------- #

def _render_pdf(doc: ExportDocument) -> bytes:
    try:
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import (
            HRFlowable,
            ListFlowable,
            ListItem,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
        )
    except ImportError as exc:  # pragma: no cover - depends on install surface
        raise ExportError(f"PDF export unavailable (reportlab not installed): {exc}") from exc

    buf = io.BytesIO()
    pdf = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=20 * mm, rightMargin=20 * mm, topMargin=18 * mm, bottomMargin=18 * mm,
        title=doc.title,
    )
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("SiftTitle", parent=styles["Title"], fontSize=18, spaceAfter=6)
    sub = ParagraphStyle("SiftSub", parent=styles["Normal"], fontSize=9, textColor=colors.grey, spaceAfter=2)
    meta = ParagraphStyle("SiftMeta", parent=styles["Normal"], fontSize=9, textColor=colors.HexColor("#444444"))
    head = ParagraphStyle("SiftHead", parent=styles["Heading2"], fontSize=12, spaceBefore=10, spaceAfter=4)
    body = ParagraphStyle("SiftBody", parent=styles["BodyText"], fontSize=10.5, leading=15, alignment=TA_LEFT)
    cite = ParagraphStyle("SiftCite", parent=styles["Normal"], fontSize=8.5, textColor=colors.HexColor("#555555"), leftIndent=6)

    def esc(text: str) -> str:
        return escape(text).replace("\n", "<br/>")

    story: list[Any] = [Paragraph(esc(doc.title), h1)]
    if doc.subtitle:
        story.append(Paragraph(esc(doc.subtitle), sub))
    for line in doc.meta_lines:
        story.append(Paragraph(esc(line), meta))
    story.append(Spacer(1, 6))
    story.append(HRFlowable(width="100%", thickness=0.6, color=colors.HexColor("#cccccc")))

    for block in doc.blocks:
        story.append(Paragraph(esc(block.heading), head))
        for para in block.paragraphs:
            story.append(Paragraph(esc(para), body))
            story.append(Spacer(1, 3))
        if block.citations:
            story.append(Paragraph("Citations", cite))
            items = [ListItem(Paragraph(esc(c), cite)) for c in block.citations]
            story.append(ListFlowable(items, bulletType="bullet", start="•"))
        story.append(Spacer(1, 4))

    pdf.build(story)
    return buf.getvalue()


def _render_docx(doc: ExportDocument) -> bytes:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError as exc:  # pragma: no cover - depends on install surface
        raise ExportError(f"DOCX export unavailable (python-docx not installed): {exc}") from exc

    document = Document()
    document.add_heading(doc.title, level=0)
    if doc.subtitle:
        p = document.add_paragraph(doc.subtitle)
        p.runs[0].italic = True
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    for line in doc.meta_lines:
        mp = document.add_paragraph(line)
        mp.runs[0].font.size = Pt(9)
        mp.runs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    for block in doc.blocks:
        document.add_heading(block.heading, level=2)
        for para in block.paragraphs:
            document.add_paragraph(para)
        if block.citations:
            cp = document.add_paragraph()
            cp.add_run("Citations").bold = True
            for c in block.citations:
                document.add_paragraph(c, style="List Bullet")

    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _render_pptx(doc: ExportDocument) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - depends on install surface
        raise ExportError(f"PPTX export unavailable (python-pptx not installed): {exc}") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide.
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = doc.title
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "\n".join(doc.meta_lines)

    blank_bullets = prs.slide_layouts[1]  # Title + Content

    # Keep slides readable: cap paragraphs per slide, continue on overflow.
    MAX_PER_SLIDE = 6
    for block in doc.blocks:
        lines = list(block.paragraphs)
        if block.citations:
            lines.append("Citations:")
            lines.extend(f"• {c}" for c in block.citations)
        if not lines:
            lines = ["(no content)"]
        for page_start in range(0, len(lines), MAX_PER_SLIDE):
            page_lines = lines[page_start : page_start + MAX_PER_SLIDE]
            slide = prs.slides.add_slide(blank_bullets)
            slide.shapes.title.text = block.heading if page_start == 0 else f"{block.heading} (cont.)"
            body_ph = slide.placeholders[1]
            tf = body_ph.text_frame
            tf.word_wrap = True
            for idx, line in enumerate(page_lines):
                para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                para.text = line if len(line) <= 500 else line[:497] + "…"
                para.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_RENDERERS = {"PDF": _render_pdf, "DOCX": _render_docx, "PPTX": _render_pptx}


def render(doc: ExportDocument, fmt: str) -> bytes:
    fmt_up = (fmt or "").upper()
    if fmt_up not in _RENDERERS:
        raise ExportError(f"Unsupported export format: {fmt!r}. Use one of {sorted(SUPPORTED_FORMATS)}.")
    return _RENDERERS[fmt_up](doc)


def export_chat(
    fmt: str,
    chat: dict[str, Any],
    messages: list[dict[str, Any]],
    *,
    generated_by: str | None = None,
    matter: dict[str, Any] | None = None,
    chambers: dict[str, Any] | None = None,
) -> tuple[bytes, str, str]:
    """Render a chat to (bytes, filename, mime_type)."""
    fmt_up = (fmt or "").upper()
    if fmt_up not in SUPPORTED_FORMATS:
        raise ExportError(f"Unsupported export format: {fmt!r}. Use one of {sorted(SUPPORTED_FORMATS)}.")
    document = build_chat_document(
        chat, messages, generated_by=generated_by, matter=matter, chambers=chambers
    )
    data = render(document, fmt_up)
    safe_title = "".join(c if c.isalnum() or c in "-_ " else "_" for c in (chat.get("title") or "memo")).strip()
    safe_title = (safe_title or "memo").replace(" ", "_")[:60]
    filename = f"{safe_title}.{EXTENSIONS[fmt_up]}"
    return data, filename, MIME_TYPES[fmt_up]
